import base64
import json
import io
from pptx import Presentation

from wasabi import msg

from goldenverba.components.document import Document, create_document
from goldenverba.components.interfaces import Reader
from goldenverba.server.types import FileConfig

# Optional imports with error handling
try:
    from pypdf import PdfReader
except ImportError:
    msg.warn("pypdf not installed, PDF functionality will be limited.")
    PdfReader = None

try:
    import docx
    import mammoth
except ImportError:
    msg.warn("python-docx or mammoth not installed, DOC/DOCX functionality will be limited.")
    docx = None
    mammoth = None

try:
    import openpyxl
    import xlrd
except ImportError:
    msg.warn("openpyxl or xlrd not installed, XLSX/XLS functionality will be limited.")
    openpyxl = None
    xlrd = None

try:
    import pandas as pd
except ImportError:
    msg.warn("pandas not installed, CSV/TSV functionality will be limited.")
    pd = None

try:
    from PIL import Image
except ImportError:
    msg.warn("Pillow (PIL) not installed, image functionality will be limited.")
    Image = None

try:
    import email
    import extract_msg
except ImportError:
    msg.warn("email or extract-msg not installed, EML/MSG functionality will be limited.")
    email = None
    extract_msg = None

try:
    import lxml.etree as ET
    from ebooklib import epub
except ImportError:
    msg.warn("lxml or ebooklib not installed, XML/EPUB functionality will be limited.")
    ET = None
    epub = None

try:
    import pyrtf
except ImportError:
    msg.warn("pyrtf not installed, RTF functionality will be limited.")
    pyrtf = None

try:
    from odf.opendocument import load as load_odt
except ImportError:
    msg.warn("odfpy not installed, ODT functionality will be limited.")
    load_odt = None

try:
    from pygments import highlight
    from pygments.lexers import get_lexer_by_name
    from pygments.formatters import HtmlFormatter
except ImportError:
    msg.warn("pygments not installed, code highlighting functionality will be limited.")
    highlight = None
    get_lexer_by_name = None
    HtmlFormatter = None

class BasicReader(Reader):
    """
    The BasicReader reads text, code, PDF, DOCX, and additional file types.
    """

    def __init__(self):
        super().__init__()
        self.name = "Default"
        self.description = "Ingests text, code, PDF, DOCX, and other common file types"
        self.requires_library = ["pypdf", "docx", "mammoth", "openpyxl", "pandas", "Pillow", "pygments", "extract_msg"]
        self.extension = [
            ".txt", ".py", ".js", ".html", ".css", ".md", ".mdx", ".json", ".pdf", ".docx", ".doc", ".pptx", ".ppt",
            ".xlsx", ".xls", ".csv", ".tsv", ".ts", ".tsx", ".vue", ".svelte", ".astro", ".php", ".rb", ".go", 
            ".rs", ".swift", ".kt", ".java", ".c", ".cpp", ".h", ".hpp", ".eml", ".msg", ".xml", ".epub", ".rtf", 
            ".odt", ".bmp", ".tiff", ".heic"
        ]  # Add supported file extensions

    async def load(self, config: dict, fileConfig: FileConfig) -> list[Document]:
        """
        Load and process a file based on its extension.
        """
        msg.info(f"Loading {fileConfig.filename} ({fileConfig.extension.lower()})")

        decoded_bytes = base64.b64decode(fileConfig.content)

        try:
            if fileConfig.extension == "":
                file_content = fileConfig.content
            elif fileConfig.extension.lower() == "json":
                return await self.load_json_file(decoded_bytes, fileConfig)
            elif fileConfig.extension.lower() == "pdf":
                file_content = await self.load_pdf_file(decoded_bytes)
            elif fileConfig.extension.lower() in ["docx", "doc"]:
                file_content = await self.load_docx_file(decoded_bytes, fileConfig.extension.lower())
            elif fileConfig.extension.lower() in [".csv", ".tsv"]:
                file_content = await self.load_csv_file(decoded_bytes, fileConfig.extension.lower())
            elif fileConfig.extension.lower() in ["xlsx", "xls"]:
                file_content = await self.load_xlsx_file(decoded_bytes, fileConfig.extension.lower())
            elif fileConfig.extension.lower() == "pptx":
                file_content = await self.load_pptx_file(decoded_bytes)
            elif fileConfig.extension.lower() in [".bmp", ".tiff", ".heic"]:
                file_content = await self.load_image_file(decoded_bytes, fileConfig.extension.lower())
            elif fileConfig.extension.lower() in [ext.lstrip(".") for ext in self.extension]:
                file_content = await self.load_text_file(decoded_bytes, fileConfig.extension.lower())
            else:
                raise ValueError(f"Unsupported file extension: {fileConfig.extension}")

            return [create_document(file_content, fileConfig)]
        except Exception as e:
            msg.fail(f"Failed to load {fileConfig.filename}: {str(e)}")
            raise

    async def load_text_file(self, decoded_bytes: bytes, extension: str) -> str:
        """Load and decode a text or code file."""
        content = decoded_bytes.decode("utf-8", errors="ignore")
        
        # Skip syntax highlighting for plain text files
        if extension == "txt":
            return content
        
        # If pygments is available and the file is not a .txt, try to apply syntax highlighting
        if highlight and get_lexer_by_name and HtmlFormatter:
            try:
                lexer = get_lexer_by_name(extension, stripall=True)
                formatter = HtmlFormatter()
                return highlight(content, lexer, formatter)
            except Exception as e:
                # If no lexer is found, return the plain content
                msg.warn(f"Syntax highlighting failed for {extension}: {str(e)}")
                return content
        
        return content


    async def load_json_file(self, decoded_bytes: bytes, fileConfig: FileConfig) -> list[Document]:
        """Load and parse a JSON file."""
        try:
            json_obj = json.loads(decoded_bytes.decode("utf-8"))
            document = Document.from_json(json_obj, None)
            return [document] if document else [create_document(json.dumps(json_obj, indent=2), fileConfig)]
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in {fileConfig.filename}: {str(e)}")

    async def load_pdf_file(self, decoded_bytes: bytes) -> str:
        """Load and extract text from a PDF file."""
        if not PdfReader:
            raise ImportError("pypdf is not installed. Cannot process PDF files.")
        pdf_bytes = io.BytesIO(decoded_bytes)
        reader = PdfReader(pdf_bytes)
        return "\n\n".join(page.extract_text() for page in reader.pages)

    async def load_docx_file(self, decoded_bytes: bytes, extension: str) -> str:
        """Load and extract text from a DOCX or DOC file."""
        if extension == "docx":
            if not docx:
                raise ImportError("python-docx is not installed. Cannot process DOCX files.")
            docx_bytes = io.BytesIO(decoded_bytes)
            doc = docx.Document(docx_bytes)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif extension == "doc":
            if not mammoth:
                raise ImportError("mammoth is not installed. Cannot process DOC files.")
            doc_bytes = io.BytesIO(decoded_bytes)
            result = mammoth.convert_to_markdown(doc_bytes)
            return result.value

    async def load_xlsx_file(self, decoded_bytes: bytes, extension: str) -> str:
        """Load and extract text from an XLSX or XLS file."""
        if extension == "xlsx":
            if not openpyxl:
                raise ImportError("openpyxl is not installed. Cannot process XLSX files.")
            xlsx_bytes = io.BytesIO(decoded_bytes)
            workbook = openpyxl.load_workbook(xlsx_bytes)
            sheet = workbook.active
            return "\n".join(["\t".join([str(cell.value) for cell in row]) for row in sheet.iter_rows()])
        elif extension == "xls":
            if not xlrd:
                raise ImportError("xlrd is not installed. Cannot process XLS files.")
            xls_bytes = io.BytesIO(decoded_bytes)
            workbook = xlrd.open_workbook(file_contents=xls_bytes.read())
            sheet = workbook.sheet_by_index(0)
            return "\n".join(["\t".join([str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]) for r in range(sheet.nrows)])

    async def load_csv_file(self, decoded_bytes: bytes, extension: str) -> str:
        """Load and extract text from a CSV/TSV file."""
        if not pd:
            raise ImportError("pandas is not installed. Cannot process CSV/TSV files.")
        csv_bytes = io.BytesIO(decoded_bytes)
        separator = "\t" if extension == ".tsv" else ","
        df = pd.read_csv(csv_bytes, sep=separator)
        return df.to_csv(index=False, sep=separator)

    async def load_pptx_file(self, decoded_bytes: bytes) -> str:
        """Load and extract text from a PPTX file."""
        if not Presentation:
            raise ImportError("python-pptx is not installed. Cannot process PPTX files.")
        pptx_bytes = io.BytesIO(decoded_bytes)
        presentation = Presentation(pptx_bytes)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n\n".join(text)

    async def load_image_file(self, decoded_bytes: bytes, extension: str) -> str:
        """Load and extract text or metadata from an image file."""
        if not Image:
            raise ImportError("Pillow (PIL) is not installed. Cannot process image files.")
        img_bytes = io.BytesIO(decoded_bytes)
        image = Image.open(img_bytes)
        # Assuming basic handling like getting image metadata or displaying the image itself
        return f"Image: {image.format}, {image.size}, {image.mode}"

    async def load_xml_file(self, decoded_bytes: bytes) -> str:
        """Load and parse an XML file."""
        if not ET:
            raise ImportError("lxml is not installed. Cannot process XML files.")
        xml_bytes = io.BytesIO(decoded_bytes)
        tree = ET.parse(xml_bytes)
        root = tree.getroot()
        return ET.tostring(root, encoding="unicode")

    async def load_epub_file(self, decoded_bytes: bytes) -> str:
        """Load and parse an EPUB file."""
        if not epub:
            raise ImportError("ebooklib is not installed. Cannot process EPUB files.")
        epub_bytes = io.BytesIO(decoded_bytes)
        book = epub.read_epub(epub_bytes)
        text = []
        for item in book.items:
            if item.get_type() == epub.EpubHtml:
                text.append(item.get_body_content().decode("utf-8"))
        return "\n\n".join(text)

    async def load_rtf_file(self, decoded_bytes: bytes) -> str:
        """Load and parse an RTF file."""
        if not pyrtf:
            raise ImportError("pyrtf is not installed. Cannot process RTF files.")
        rtf_bytes = io.BytesIO(decoded_bytes)
        doc = pyrtf.parse(rtf_bytes)
        return doc.as_text()

    async def load_odt_file(self, decoded_bytes: bytes) -> str:
        """Load and parse an ODT file."""
        if not load_odt:
            raise ImportError("odfpy is not installed. Cannot process ODT files.")
        odt_bytes = io.BytesIO(decoded_bytes)
        doc = load_odt(odt_bytes)
        # Return the text of the ODT file (assuming simple text extraction)
        return "\n".join(paragraph.text for paragraph in doc.text)