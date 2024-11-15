import base64
import io
import os
import subprocess
import tempfile
from typing import List, Optional, Tuple
from wasabi import msg
from goldenverba.components.document import Document, create_document
from goldenverba.components.interfaces import Reader
from goldenverba.server.types import FileConfig

# Primary unstructured imports with error handling
try:
    from unstructured.partition.auto import partition
    from unstructured.partition.pdf import partition_pdf
    from unstructured.partition.docx import partition_docx
    from unstructured.partition.doc import partition_doc
    from unstructured.partition.pptx import partition_pptx
    from unstructured.partition.ppt import partition_ppt
    from unstructured.partition.xlsx import partition_xlsx
    from unstructured.partition.csv import partition_csv
    from unstructured.partition.text import partition_text
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    msg.warn("unstructured library not installed, falling back to alternative methods")
    UNSTRUCTURED_AVAILABLE = False

# Fallback imports - only import if needed
PyPDF2 = None
pdfplumber = None
docx = None
Presentation = None

class CustomUnstructuredReader(Reader):
    """
    A reader that uses the Unstructured library primarily, with LibreOffice conversion
    for legacy formats and additional fallbacks when necessary.
    Implements the Reader interface from goldenverba.components.interfaces.
    """

    def __init__(self):
        super().__init__()
        self.name = "UnstructuredLibrary"
        self.description = "Ingests documents using Unstructured library with LibreOffice conversion and fallbacks"
        self.requires_library = ["unstructured"]
        self.extension = [
            "pdf", "xlsx", "csv", "docx", "doc", "pptx", "ppt", "txt"
        ]
        self._check_dependencies()

    def _check_dependencies(self) -> None:
        """Check system dependencies for unstructured partitioning and LibreOffice."""
        self.has_poppler = False
        self.has_libreoffice = False
        
        try:
            subprocess.run(['pdfinfo', '-v'], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            self.has_poppler = True
        except FileNotFoundError:
            msg.warn("Poppler not found - PDF processing may be limited")

        try:
            process = subprocess.run(['soffice', '--version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            self.has_libreoffice = True
            msg.good(f"LibreOffice found: {process.stdout.decode().strip()}")
        except FileNotFoundError:
            msg.warn("LibreOffice not found - DOC/PPT processing will be limited")

    def _load_fallback_libraries(self, file_type: str) -> None:
        """Lazily load fallback libraries only when needed."""
        global PyPDF2, pdfplumber, docx, Presentation

        if file_type == "pdf" and not (PyPDF2 and pdfplumber):
            try:
                import PyPDF2
            except ImportError:
                msg.warn("PyPDF2 not installed")
            try:
                import pdfplumber
            except ImportError:
                msg.warn("pdfplumber not installed")

        elif file_type in ["doc", "docx"] and not docx:
            try:
                import docx
            except ImportError:
                msg.warn("python-docx not installed")

        elif file_type in ["ppt", "pptx"] and not Presentation:
            try:
                from pptx import Presentation
            except ImportError:
                msg.warn("python-pptx not installed")

    def _convert_using_libreoffice(self, file_bytes: bytes, input_ext: str, output_ext: str) -> Optional[bytes]:
        """
        Convert a document using LibreOffice.
        Args:
            file_bytes: Input file content
            input_ext: Input file extension (e.g., 'ppt', 'doc')
            output_ext: Desired output extension (e.g., 'pptx', 'docx')
        Returns:
            Converted file bytes or None if conversion fails
        """
        if not self.has_libreoffice:
            return None

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                # Write input file
                input_path = os.path.join(temp_dir, f"input.{input_ext}")
                with open(input_path, "wb") as f:
                    f.write(file_bytes)

                # Create output directory
                output_dir = os.path.join(temp_dir, "output")
                os.makedirs(output_dir, exist_ok=True)
                
                # Run LibreOffice conversion
                process = subprocess.run([
                    'soffice',
                    '--headless',
                    '--convert-to', output_ext,
                    '--outdir', output_dir,
                    input_path
                ], capture_output=True)

                if process.returncode != 0:
                    msg.warn(f"LibreOffice conversion failed: {process.stderr.decode()}")
                    return None

                # Read converted file
                output_path = os.path.join(output_dir, f"input.{output_ext}")
                if os.path.exists(output_path):
                    with open(output_path, "rb") as f:
                        return f.read()

        except Exception as e:
            msg.warn(f"LibreOffice conversion failed: {str(e)}")

        return None

    async def _process_pdf(self, file_bytes: bytes) -> str:
        """Process PDF with unstructured first, then fallbacks if needed."""
        try:
            if self.has_poppler:
                file_obj = io.BytesIO(file_bytes)
                elements = partition_pdf(file=file_obj, strategy="fast")
                text = "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
                if text.strip():
                    return text
                msg.warn("Unstructured PDF processing produced no text, trying fallbacks")
        except Exception as e:
            msg.warn(f"Unstructured PDF processing failed: {str(e)}")

        # Load fallback libraries only if needed
        self._load_fallback_libraries("pdf")
        
        # Try PyPDF2
        if PyPDF2:
            try:
                with io.BytesIO(file_bytes) as pdf_file:
                    reader = PyPDF2.PdfReader(pdf_file)
                    text = "\n\n".join(page.extract_text() for page in reader.pages)
                    if text.strip():
                        return text
            except Exception as e:
                msg.warn(f"PyPDF2 extraction failed: {str(e)}")

        # Try pdfplumber as last resort
        if pdfplumber:
            try:
                with io.BytesIO(file_bytes) as pdf_file:
                    with pdfplumber.open(pdf_file) as pdf:
                        text = "\n\n".join(page.extract_text() for page in pdf.pages)
                        if text.strip():
                            return text
            except Exception as e:
                msg.warn(f"pdfplumber extraction failed: {str(e)}")

        msg.warn("All PDF extraction methods failed")
        return ""

    async def _process_docx(self, file_bytes: bytes) -> str:
        """Process DOCX with unstructured first, then fallback if needed."""
        try:
            file_obj = io.BytesIO(file_bytes)
            elements = partition_docx(file=file_obj)
            text = "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
            if text.strip():
                return text
            msg.warn("Unstructured DOCX processing produced no text, trying fallback")
        except Exception as e:
            msg.warn(f"Unstructured DOCX processing failed: {str(e)}")

        # Load fallback library only if needed
        self._load_fallback_libraries("docx")
        
        if docx:
            try:
                doc = docx.Document(io.BytesIO(file_bytes))
                return "\n".join(paragraph.text for paragraph in doc.paragraphs)
            except Exception as e:
                msg.warn(f"DOCX fallback extraction failed: {str(e)}")
        
        return ""

    async def _process_doc(self, file_bytes: bytes) -> str:
        """Process DOC by converting to DOCX first, then fallback to other methods."""
        # Try converting to DOCX first
        if self.has_libreoffice:
            msg.info("Attempting to convert DOC to DOCX using LibreOffice")
            docx_bytes = self._convert_using_libreoffice(file_bytes, "doc", "docx")
            if docx_bytes:
                msg.good("Successfully converted DOC to DOCX")
                return await self._process_docx(docx_bytes)
            msg.warn("DOC to DOCX conversion failed, trying direct processing")

        # If conversion fails, try direct processing
        try:
            if self.has_libreoffice:
                file_obj = io.BytesIO(file_bytes)
                elements = partition_doc(file=file_obj)
                text = "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
                if text.strip():
                    return text
        except Exception as e:
            msg.warn(f"Direct DOC processing failed: {str(e)}")

        # Last resort: try basic text extraction
        try:
            text = file_bytes.decode('utf-8', errors='ignore')
            if text.strip():
                return text
        except Exception as e:
            msg.warn(f"Basic text extraction failed: {str(e)}")
        
        return ""

    async def _process_pptx(self, file_bytes: bytes) -> str:
        """Process PPTX with unstructured first, then fallback if needed."""
        try:
            file_obj = io.BytesIO(file_bytes)
            elements = partition_pptx(file=file_obj)
            text = "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
            if text.strip():
                return text
            msg.warn("Unstructured PPTX processing produced no text, trying fallback")
        except Exception as e:
            msg.warn(f"Unstructured PPTX processing failed: {str(e)}")

        # Load fallback library only if needed
        self._load_fallback_libraries("pptx")
        
        if Presentation:
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return "\n\n".join(text_runs)
            except Exception as e:
                msg.warn(f"PPTX fallback extraction failed: {str(e)}")
        
        return ""

    async def _process_ppt(self, file_bytes: bytes) -> str:
        """Process PPT by converting to PPTX first, then fallback to other methods."""
        # Try converting to PPTX first
        if self.has_libreoffice:
            msg.info("Attempting to convert PPT to PPTX using LibreOffice")
            pptx_bytes = self._convert_using_libreoffice(file_bytes, "ppt", "pptx")
            if pptx_bytes:
                msg.good("Successfully converted PPT to PPTX")
                return await self._process_pptx(pptx_bytes)
            msg.warn("PPT to PPTX conversion failed, trying direct processing")

        # If conversion fails, try direct processing
        try:
            if self.has_libreoffice:
                file_obj = io.BytesIO(file_bytes)
                elements = partition_ppt(file=file_obj)
                text = "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
                if text.strip():
                    return text
        except Exception as e:
            msg.warn(f"Direct PPT processing failed: {str(e)}")

        # Last resort: try basic text extraction
        try:
            text = file_bytes.decode('utf-16le', errors='ignore')
            if text.strip():
                return text
        except Exception as e:
            try:
                text = file_bytes.decode('utf-8', errors='ignore')
                if text.strip():
                    return text
            except Exception as e:
                msg.warn(f"Basic text extraction failed: {str(e)}")
        
        return ""

    async def _process_xlsx(self, file_bytes: bytes) -> str:
        """Process XLSX with unstructured."""
        try:
            file_obj = io.BytesIO(file_bytes)
            elements = partition_xlsx(file=file_obj)
            return "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
        except Exception as e:
            msg.warn(f"XLSX processing failed: {str(e)}")
            return ""

    async def _process_csv(self, file_bytes: bytes) -> str:
        """Process CSV with unstructured."""
        try:
            file_obj = io.BytesIO(file_bytes)
            elements = partition_csv(file=file_obj)
            return "\n\n".join([str(el) for el in elements if hasattr(el, 'text')])
        except Exception as e:
            msg.warn(f"CSV processing failed: {str(e)}")
            return ""

    async def load(self, config: dict, fileConfig: FileConfig) -> List[Document]:
        """Load and process a file, returning a list of Document objects."""
        try:
            msg.info(f"Processing {fileConfig.filename}")
            decoded_bytes = base64.b64decode(fileConfig.content)
            extension = fileConfig.extension.lower()

            # Process based on file extension
            content = ""
            if extension == "pdf":
                content = await self._process_pdf(decoded_bytes)
            elif extension == "docx":
                content = await self._process_docx(decoded_bytes)
            elif extension == "doc":
                content = await self._process_doc(decoded_bytes)
            elif extension == "pptx":
                content = await self._process_pptx(decoded_bytes)
            elif extension == "ppt":
                content = await self._process_ppt(decoded_bytes)
            elif extension == "xlsx":
                content = await self._process_xlsx(decoded_bytes)
            elif extension == "csv":
                content = await self._process_csv(decoded_bytes)
            elif extension == "txt":
                content = decoded_bytes.decode('utf-8', errors='ignore')
            else:
                msg.warn(f"Unsupported extension: {extension}")
                content = ""

            return [create_document(content, fileConfig)]

        except Exception as e:
            msg.fail(f"Failed to process {fileConfig.filename}: {str(e)}")
            raise

    def __str__(self) -> str:
        return f"{self.name} - {self.description}"

    def __repr__(self) -> str:
        return f"CustomUnstructuredReader(extensions={self.extension})"
